"""WEB3 daily news -> PowerPoint generator.

Parses docs/YYYY-MM-DD.md, then builds one densely packed slide per date.
Heavy on shapes, color-coded categories, and call-out boxes so the day's
big stories stand out at a glance.
"""

from __future__ import annotations

import glob
import os
import re
from collections import Counter
from dataclasses import dataclass, field
from typing import List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ---------- color palette (dark cyber theme) ----------
BG          = RGBColor(0x0B, 0x14, 0x2A)
PANEL       = RGBColor(0x12, 0x1E, 0x3D)
PANEL_SOFT  = RGBColor(0x1A, 0x29, 0x4F)
GRID_LINE   = RGBColor(0x27, 0x3A, 0x66)
TEXT        = RGBColor(0xE6, 0xED, 0xF7)
TEXT_DIM    = RGBColor(0x9A, 0xA8, 0xC4)
ACCENT      = RGBColor(0x00, 0xE5, 0xFF)   # cyan
HOT         = RGBColor(0xFF, 0x3D, 0x8A)   # magenta - 重大
WARN        = RGBColor(0xFF, 0xC1, 0x07)   # amber - 注目
GOOD        = RGBColor(0x4A, 0xE3, 0xA3)   # mint - ポジティブ
COOL        = RGBColor(0x6E, 0x8B, 0xFF)   # periwinkle

# category color map
CATEGORY_COLORS = {
    "規制・政策":     RGBColor(0xFF, 0xC1, 0x07),
    "ETF・機関":      RGBColor(0x4A, 0xE3, 0xA3),
    "価格・市場":     RGBColor(0xFF, 0x6B, 0x3D),
    "ステーブルコイン": RGBColor(0x6E, 0x8B, 0xFF),
    "DeFi・L2":       RGBColor(0x9B, 0x59, 0xE6),
    "セキュリティ":   RGBColor(0xFF, 0x3D, 0x8A),
    "国内企業":       RGBColor(0x00, 0xE5, 0xFF),
    "NFT・ゲーム":    RGBColor(0xF4, 0x7A, 0xC9),
    "AI・技術":       RGBColor(0x4B, 0xC8, 0xF0),
    "マイニング":     RGBColor(0xB0, 0x88, 0x55),
    "その他":         RGBColor(0x9A, 0xA8, 0xC4),
}

CATEGORY_KEYWORDS = [
    ("規制・政策",     ["SEC", "CFTC", "規制", "法案", "CLARITY", "GENIUS", "ホワイトハウス",
                       "上院", "下院", "議会", "政府", "FRB", "FOMC", "課税", "税制",
                       "Trump", "トランプ", "金融庁", "EU", "MiCA", "禁止", "承認", "免許"]),
    ("ETF・機関",      ["ETF", "ETP", "BlackRock", "ブラックロック", "Fidelity", "Grayscale",
                       "機関投資家", "上場", "IPO", "Strategy", "ストラテジー", "MicroStrategy"]),
    ("価格・市場",     ["価格", "急騰", "急落", "暴落", "ATH", "高値", "安値", "下落", "上昇",
                       "相場", "市場", "ボラ", "清算", "ロング", "ショート", "オプション",
                       "先物", "出来高", "デリバ"]),
    ("ステーブルコイン", ["ステーブル", "USDT", "USDC", "JPYC", "PYUSD", "Tether", "Circle",
                          "RLUSD", "stablecoin"]),
    ("DeFi・L2",       ["DeFi", "L2", "レイヤー", "Layer", "DEX", "Uniswap", "Aave",
                       "Optimism", "Arbitrum", "Base", "Polygon", "zkSync", "Scroll",
                       "ロールアップ", "ハイパーリキッド", "Hyperliquid"]),
    ("セキュリティ",   ["ハッキング", "ハック", "流出", "exploit", "脆弱", "盗難",
                       "詐欺", "scam", "rug", "攻撃", "セキュリティ", "凍結"]),
    ("国内企業",       ["JPYC", "リミックス", "メタプラ", "Metaplanet", "SBI", "GMO",
                       "bitFlyer", "Coincheck", "三菱UFJ", "Progmat", "野村", "MUFG",
                       "ANA", "JAL", "東証", "日本", "国内"]),
    ("NFT・ゲーム",    ["NFT", "ゲーム", "GameFi", "メタバース", "OpenSea", "Magic Eden",
                       "Pudgy", "BAYC", "Web3ゲーム"]),
    ("AI・技術",       ["AI", "人工知能", "OpenAI", "ChatGPT", "Claude", "エージェント",
                       "RWA", "トークン化", "ZK", "ゼロ知識"]),
    ("マイニング",     ["マイニング", "マイナー", "ハッシュレート", "Mining", "miner",
                       "Bitcoin Depot", "ATM"]),
]

# weight terms used to flag "今日の重大ニュース"
HOT_TERMS = [
    "緊急", "速報", "破産", "破綻", "倒産", "ハッキング", "ハック", "流出",
    "崩壊", "暴落", "急落", "ATH", "過去最高", "史上最高",
    "可決", "成立", "承認", "施行", "発効",
    "提訴", "起訴", "逮捕",
    "数十億", "1兆", "兆円", "億ドル",
    "CLARITY", "GENIUS", "FOMC", "FRB",
    "Strategy", "BlackRock", "Trump", "トランプ", "ホワイトハウス",
]

MEDIA_TAG = re.compile(r"\*\*メディア\*\*:\s*([^\n]+)")
TIME_TAG  = re.compile(r"\*\*公開日時\*\*:\s*([^\n]+)")
SUMMARY_TAG = re.compile(r"\*\*概要（原文）\*\*:\s*([^\n]+)")
ART_HEAD  = re.compile(r"^###\s+\d+\.\s+\[(?P<title>.+?)\]\((?P<url>.+?)\)\s*$")


@dataclass
class Article:
    idx: int
    title: str
    url: str
    media: str = ""
    published: str = ""
    summary: str = ""
    section: str = ""   # 国内 / 海外
    category: str = "その他"
    hot_score: int = 0

    @property
    def hhmm(self) -> str:
        m = re.search(r"(\d{2}:\d{2})", self.published)
        return m.group(1) if m else ""


@dataclass
class DayReport:
    date: str
    domestic_count: int = 0
    overseas_count: int = 0
    articles: List[Article] = field(default_factory=list)


# ---------- parsing ----------
def parse_md(path: str) -> DayReport:
    text = open(path, "r", encoding="utf-8").read()
    date = os.path.basename(path).replace(".md", "")
    rep = DayReport(date=date)

    m = re.search(r"国内:\s*(\d+)件\s*/\s*海外:\s*(\d+)件", text)
    if m:
        rep.domestic_count = int(m.group(1))
        rep.overseas_count = int(m.group(2))

    section = ""
    lines = text.splitlines()
    cur: Article | None = None
    for line in lines:
        if line.startswith("## "):
            if "国内" in line:
                section = "国内"
            elif "海外" in line:
                section = "海外"
            else:
                section = ""
            continue
        head = ART_HEAD.match(line)
        if head:
            if cur:
                cur.category = classify(cur)
                cur.hot_score = hot_score(cur)
                rep.articles.append(cur)
            idx = int(re.match(r"^###\s+(\d+)\.", line).group(1))
            cur = Article(idx=idx, title=head.group("title"), url=head.group("url"),
                          section=section)
            continue
        if cur is None:
            continue
        if (mm := MEDIA_TAG.search(line)):
            cur.media = mm.group(1).strip()
        elif (mm := TIME_TAG.search(line)):
            cur.published = mm.group(1).strip()
        elif (mm := SUMMARY_TAG.search(line)):
            cur.summary = mm.group(1).strip()
    if cur:
        cur.category = classify(cur)
        cur.hot_score = hot_score(cur)
        rep.articles.append(cur)
    return rep


def classify(a: Article) -> str:
    blob = a.title + " " + a.summary
    for cat, kws in CATEGORY_KEYWORDS:
        for kw in kws:
            if kw.lower() in blob.lower():
                return cat
    return "その他"


def hot_score(a: Article) -> int:
    blob = a.title + " " + a.summary
    score = 0
    for kw in HOT_TERMS:
        if kw.lower() in blob.lower():
            score += 2
    # any cash-amount mention is a small bump
    if re.search(r"\d+(億|兆|百万ドル|億ドル|兆ドル)", blob):
        score += 1
    if a.section == "国内":
        score += 0  # neutral; domestic already smaller pool
    return score


# ---------- pptx helpers ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def set_line(shape, color: RGBColor, width: float = 0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_rect(slide, x, y, w, h, fill=PANEL, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    set_fill(s, fill)
    if line is None:
        set_no_line(s)
    else:
        set_line(s, line, 0.75)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=10, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Yu Gothic UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb


def add_multiline(slide, x, y, w, h, runs: List[Tuple[str, dict]],
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, style-dict). style keys: size, color, bold, font, line_break_after."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, style in runs:
        if style.get("new_paragraph") and not first:
            p = tf.add_paragraph()
            p.alignment = style.get("align", align)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(style.get("size", 10))
        r.font.bold = style.get("bold", False)
        r.font.name = style.get("font", "Yu Gothic UI")
        r.font.color.rgb = style.get("color", TEXT)
        first = False
    return tb


# ---------- slide builder ----------
def build_index_slide(prs: Presentation, reports: List[DayReport]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)

    # title bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill=PANEL)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
             "WEB3 NEWS DIGEST  /  2026-05-07 → 2026-05-18",
             size=28, color=ACCENT, bold=True)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
             "暗号資産・ブロックチェーン業界 12日間の総まとめ ─ デイリーレポート",
             size=14, color=TEXT_DIM)

    # totals
    total_d = sum(r.domestic_count for r in reports)
    total_o = sum(r.overseas_count for r in reports)
    total = total_d + total_o
    cards = [
        ("収録日数", f"{len(reports)} 日", ACCENT),
        ("総記事数", f"{total:,} 本", GOOD),
        ("国内", f"{total_d:,} 本", COOL),
        ("海外", f"{total_o:,} 本", WARN),
    ]
    cw = Inches(2.9); gap = Inches(0.2); cx = Inches(0.5)
    cy = Inches(1.35)
    for label, val, col in cards:
        add_rect(slide, cx, cy, cw, Inches(1.05), fill=PANEL, line=GRID_LINE)
        add_rect(slide, cx, cy, Inches(0.12), Inches(1.05), fill=col)
        add_text(slide, cx + Inches(0.25), cy + Inches(0.08), cw, Inches(0.35),
                 label, size=11, color=TEXT_DIM)
        add_text(slide, cx + Inches(0.25), cy + Inches(0.40), cw, Inches(0.65),
                 val, size=24, color=col, bold=True)
        cx += cw + gap

    # daily bar chart of totals
    chart_x = Inches(0.5); chart_y = Inches(2.65)
    chart_w = Inches(12.3); chart_h = Inches(2.2)
    add_rect(slide, chart_x, chart_y, chart_w, chart_h, fill=PANEL, line=GRID_LINE)
    add_text(slide, chart_x + Inches(0.2), chart_y + Inches(0.05),
             chart_w, Inches(0.35),
             "日別ニュース本数（国内 + 海外）", size=12, color=ACCENT, bold=True)
    max_total = max((r.domestic_count + r.overseas_count) for r in reports)
    bar_area_x = chart_x + Inches(0.4)
    bar_area_y = chart_y + Inches(0.55)
    bar_area_w = chart_w - Inches(0.8)
    bar_area_h = chart_h - Inches(1.0)
    n = len(reports)
    slot = bar_area_w / n
    bar_w = slot * 0.65
    for i, r in enumerate(reports):
        tot = r.domestic_count + r.overseas_count
        dom = r.domestic_count
        ov = r.overseas_count
        bx = bar_area_x + slot * i + (slot - bar_w) / 2
        full_h = bar_area_h * (tot / max_total)
        dom_h = bar_area_h * (dom / max_total)
        ov_h = full_h - dom_h
        # overseas (bottom of stack visually = upper part of bar drawn from top down)
        ov_y = bar_area_y + bar_area_h - full_h
        b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, ov_y, bar_w, ov_h)
        set_fill(b1, WARN); set_no_line(b1); b1.shadow.inherit = False
        b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, ov_y + ov_h, bar_w, dom_h)
        set_fill(b2, COOL); set_no_line(b2); b2.shadow.inherit = False
        # value
        add_text(slide, bx - Inches(0.1), ov_y - Inches(0.3), bar_w + Inches(0.2), Inches(0.28),
                 f"{tot}", size=8, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        # date label
        add_text(slide, bx - Inches(0.1), bar_area_y + bar_area_h + Inches(0.05),
                 bar_w + Inches(0.2), Inches(0.3),
                 r.date[5:], size=8, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # legend
    lx = chart_x + chart_w - Inches(2.6)
    ly = chart_y + Inches(0.08)
    add_rect(slide, lx, ly + Inches(0.05), Inches(0.18), Inches(0.18), fill=COOL)
    add_text(slide, lx + Inches(0.25), ly, Inches(1.1), Inches(0.3),
             "国内", size=10, color=TEXT)
    add_rect(slide, lx + Inches(1.0), ly + Inches(0.05), Inches(0.18), Inches(0.18), fill=WARN)
    add_text(slide, lx + Inches(1.25), ly, Inches(1.1), Inches(0.3),
             "海外", size=10, color=TEXT)

    # cross-period themes (top categories overall)
    all_articles: List[Article] = []
    for r in reports:
        all_articles.extend(r.articles)
    cat_counts = Counter(a.category for a in all_articles)
    cat_top = cat_counts.most_common()

    box_x = Inches(0.5); box_y = Inches(4.95)
    box_w = Inches(6.2); box_h = Inches(2.35)
    add_rect(slide, box_x, box_y, box_w, box_h, fill=PANEL, line=GRID_LINE)
    add_text(slide, box_x + Inches(0.2), box_y + Inches(0.08),
             box_w, Inches(0.35),
             "期間内の主要テーマ（記事数）", size=12, color=ACCENT, bold=True)
    # horizontal bars
    max_cat = max(c for _, c in cat_top) if cat_top else 1
    row_y = box_y + Inches(0.55)
    row_h = Inches(0.22)
    row_gap = Inches(0.05)
    for cat, cnt in cat_top[:7]:
        col = CATEGORY_COLORS.get(cat, TEXT_DIM)
        add_text(slide, box_x + Inches(0.2), row_y, Inches(1.8), row_h,
                 cat, size=10, color=TEXT)
        bar_max_w = Inches(3.6)
        bw = bar_max_w * (cnt / max_cat)
        bx = box_x + Inches(2.1)
        add_rect(slide, bx, row_y + Inches(0.04), bar_max_w, Inches(0.14),
                 fill=PANEL_SOFT)
        add_rect(slide, bx, row_y + Inches(0.04), bw, Inches(0.14), fill=col)
        add_text(slide, bx + bar_max_w + Inches(0.05), row_y, Inches(0.8), row_h,
                 f"{cnt}", size=10, color=TEXT, bold=True)
        row_y += row_h + row_gap

    # navigation hint
    nav_x = Inches(7.0); nav_y = Inches(4.95)
    nav_w = Inches(5.83); nav_h = Inches(2.35)
    add_rect(slide, nav_x, nav_y, nav_w, nav_h, fill=PANEL, line=GRID_LINE)
    add_text(slide, nav_x + Inches(0.2), nav_y + Inches(0.08),
             nav_w, Inches(0.35),
             "本資料の構成", size=12, color=ACCENT, bold=True)
    intro = [
        ("• 1スライド = 1日 のデイリーダイジェスト構成です。", {"size": 11, "color": TEXT}),
        ("• 各スライドは「重大ニュース」「カテゴリ別ハイライト」「主要記事グリッド」「メディア統計」で構成。",
         {"size": 11, "color": TEXT, "new_paragraph": True}),
        ("• 重大度はキーワード分析（規制成立/破綻/価格急変/巨額流出 等）で自動採点しています。",
         {"size": 11, "color": TEXT, "new_paragraph": True}),
        ("• カテゴリ色分け：規制=黄 / ETF=緑 / 価格=橙 / ステーブル=青 / DeFi=紫 / セキュリティ=ピンク / 国内=シアン。",
         {"size": 11, "color": TEXT_DIM, "new_paragraph": True}),
        (f"• 全 {len(reports)} 日分のスライドが次ページ以降に続きます（{reports[0].date} → {reports[-1].date}）。",
         {"size": 11, "color": TEXT_DIM, "new_paragraph": True}),
    ]
    add_multiline(slide, nav_x + Inches(0.2), nav_y + Inches(0.5),
                  nav_w - Inches(0.4), nav_h - Inches(0.6), intro)


def build_day_slide(prs: Presentation, rep: DayReport):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)

    # ============ Header ============
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), fill=PANEL)
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Inches(0.04), fill=ACCENT)
    add_text(slide, Inches(0.3), Inches(0.10), Inches(5.5), Inches(0.45),
             f"DAILY DIGEST  /  {rep.date}",
             size=22, color=ACCENT, bold=True)
    weekday = ["月", "火", "水", "木", "金", "土", "日"]
    try:
        import datetime as _dt
        d = _dt.date.fromisoformat(rep.date)
        wd = weekday[d.weekday()]
    except Exception:
        wd = ""
    add_text(slide, Inches(0.3), Inches(0.50), Inches(5.5), Inches(0.35),
             f"WEB3 / Crypto News Roundup ({wd})", size=11, color=TEXT_DIM)

    # KPI mini cards in header (right side)
    total = rep.domestic_count + rep.overseas_count
    pos = Inches(5.9)
    for label, val, col in [
        ("総数", f"{total}", ACCENT),
        ("国内", f"{rep.domestic_count}", COOL),
        ("海外", f"{rep.overseas_count}", WARN),
        ("重大", f"{sum(1 for a in rep.articles if a.hot_score >= 4)}", HOT),
    ]:
        add_rect(slide, pos, Inches(0.12), Inches(1.7), Inches(0.62),
                 fill=PANEL_SOFT, line=GRID_LINE)
        add_text(slide, pos + Inches(0.1), Inches(0.14), Inches(1.5), Inches(0.25),
                 label, size=9, color=TEXT_DIM)
        add_text(slide, pos + Inches(0.1), Inches(0.34), Inches(1.5), Inches(0.4),
                 val, size=18, color=col, bold=True)
        pos += Inches(1.78)

    # ============ Left column: TOP HEADLINES ============
    L_X = Inches(0.25); L_Y = Inches(1.05); L_W = Inches(7.45); L_H = Inches(6.30)
    add_rect(slide, L_X, L_Y, L_W, L_H, fill=PANEL, line=GRID_LINE)
    add_rect(slide, L_X, L_Y, Inches(0.12), L_H, fill=HOT)
    add_text(slide, L_X + Inches(0.25), L_Y + Inches(0.08), L_W, Inches(0.35),
             "■ 今日の重大ニュース TOP 5", size=13, color=HOT, bold=True)
    add_text(slide, L_X + Inches(0.25), L_Y + Inches(0.40), L_W, Inches(0.3),
             "規制成立 / 破綻 / 価格急変 / 巨額資金移動 等を自動抽出",
             size=9, color=TEXT_DIM)

    # rank articles
    ranked = sorted(rep.articles, key=lambda a: (-a.hot_score, a.idx))[:5]
    # ensure we have 5 even if scores are zero
    if len(ranked) < 5:
        rest = [a for a in rep.articles if a not in ranked]
        ranked.extend(rest[: 5 - len(ranked)])

    card_y = L_Y + Inches(0.78)
    card_h = Inches(1.05)
    card_gap = Inches(0.05)
    rank_colors = [HOT, HOT, WARN, WARN, GOOD]
    for i, a in enumerate(ranked):
        col = rank_colors[i] if i < len(rank_colors) else COOL
        cy = card_y + i * (card_h + card_gap)
        cx = L_X + Inches(0.2)
        cw = L_W - Inches(0.4)
        add_rect(slide, cx, cy, cw, card_h, fill=PANEL_SOFT, line=GRID_LINE)
        # rank badge
        add_rect(slide, cx, cy, Inches(0.55), card_h, fill=col)
        add_text(slide, cx, cy + Inches(0.18), Inches(0.55), Inches(0.4),
                 f"#{i+1}", size=20, color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cx, cy + Inches(0.62), Inches(0.55), Inches(0.3),
                 a.section, size=8, color=BG, bold=True, align=PP_ALIGN.CENTER)
        # body
        tx = cx + Inches(0.7); tw = cw - Inches(0.8)
        title = a.title
        if len(title) > 75:
            title = title[:73] + "…"
        meta = f"{a.media}  ·  {a.hhmm}  ·  {a.category}"
        summary = (a.summary or "").strip()
        if len(summary) > 130:
            summary = summary[:128] + "…"
        runs = [
            (title, {"size": 12, "bold": True, "color": TEXT}),
            (meta, {"size": 8.5, "color": col, "new_paragraph": True, "bold": True}),
            (summary, {"size": 9, "color": TEXT_DIM, "new_paragraph": True}),
        ]
        add_multiline(slide, tx, cy + Inches(0.06), tw, card_h - Inches(0.12), runs)

    # ============ Right column ============
    R_X = Inches(7.80); R_Y = Inches(1.05); R_W = Inches(5.28); R_H = Inches(6.30)

    # ----- Right top: category donut + counts -----
    RT_H = Inches(2.40)
    add_rect(slide, R_X, R_Y, R_W, RT_H, fill=PANEL, line=GRID_LINE)
    add_rect(slide, R_X, R_Y, Inches(0.10), RT_H, fill=ACCENT)
    add_text(slide, R_X + Inches(0.2), R_Y + Inches(0.05), R_W, Inches(0.35),
             "■ カテゴリ別 分布", size=12, color=ACCENT, bold=True)

    cat_counts = Counter(a.category for a in rep.articles)
    top_cats = cat_counts.most_common(7)
    # left: vertical color-coded list
    list_x = R_X + Inches(0.2)
    list_y = R_Y + Inches(0.45)
    list_w = Inches(2.5)
    max_cat = max(c for _, c in top_cats) if top_cats else 1
    rh = Inches(0.23)
    for cat, cnt in top_cats:
        col = CATEGORY_COLORS.get(cat, TEXT_DIM)
        add_rect(slide, list_x, list_y + Inches(0.06), Inches(0.12), Inches(0.12), fill=col)
        add_text(slide, list_x + Inches(0.18), list_y, Inches(1.6), rh,
                 cat, size=9, color=TEXT)
        add_text(slide, list_x + Inches(1.7), list_y, Inches(0.6), rh,
                 f"{cnt}", size=9, color=col, bold=True, align=PP_ALIGN.RIGHT)
        list_y += rh

    # right: stacked horizontal bar showing category mix
    bar_x = R_X + Inches(2.85)
    bar_y = R_Y + Inches(0.55)
    bar_w = R_W - Inches(3.1)
    bar_h = Inches(1.55)
    add_text(slide, bar_x, R_Y + Inches(0.30), bar_w, Inches(0.22),
             "比率", size=9, color=TEXT_DIM)
    total_arts = sum(cnt for _, cnt in cat_counts.items())
    if total_arts > 0:
        cy = bar_y
        for cat, cnt in cat_counts.most_common():
            col = CATEGORY_COLORS.get(cat, TEXT_DIM)
            seg_h = bar_h * (cnt / total_arts)
            if seg_h < Inches(0.04):
                continue
            add_rect(slide, bar_x, cy, bar_w, seg_h, fill=col)
            # label inside if room
            if seg_h > Inches(0.18):
                pct = cnt / total_arts * 100
                add_text(slide, bar_x + Inches(0.05), cy + Inches(0.02),
                         bar_w - Inches(0.1), seg_h - Inches(0.04),
                         f"{cat} {pct:.0f}%", size=8, color=BG, bold=True)
            cy += seg_h

    # ----- Right middle: media leaderboard -----
    RM_Y = R_Y + RT_H + Inches(0.07)
    RM_H = Inches(1.65)
    add_rect(slide, R_X, RM_Y, R_W, RM_H, fill=PANEL, line=GRID_LINE)
    add_rect(slide, R_X, RM_Y, Inches(0.10), RM_H, fill=WARN)
    add_text(slide, R_X + Inches(0.2), RM_Y + Inches(0.05), R_W, Inches(0.3),
             "■ メディア発信ランキング", size=12, color=WARN, bold=True)
    media_counts = Counter(a.media for a in rep.articles if a.media).most_common(6)
    if media_counts:
        max_m = media_counts[0][1]
        my = RM_Y + Inches(0.42)
        rh2 = Inches(0.18)
        for media, cnt in media_counts:
            add_text(slide, R_X + Inches(0.25), my, Inches(2.4), rh2,
                     media[:20], size=9, color=TEXT)
            bar_max = Inches(2.0)
            bw = bar_max * (cnt / max_m)
            bx2 = R_X + Inches(2.7)
            add_rect(slide, bx2, my + Inches(0.04), bar_max, Inches(0.1), fill=PANEL_SOFT)
            add_rect(slide, bx2, my + Inches(0.04), bw, Inches(0.1), fill=WARN)
            add_text(slide, bx2 + bar_max + Inches(0.05), my, Inches(0.4), rh2,
                     f"{cnt}", size=9, color=WARN, bold=True)
            my += rh2

    # ----- Right bottom: trend keyword cloud -----
    RB_Y = RM_Y + RM_H + Inches(0.07)
    RB_H = R_Y + R_H - RB_Y
    add_rect(slide, R_X, RB_Y, R_W, RB_H, fill=PANEL, line=GRID_LINE)
    add_rect(slide, R_X, RB_Y, Inches(0.10), RB_H, fill=GOOD)
    add_text(slide, R_X + Inches(0.2), RB_Y + Inches(0.05), R_W, Inches(0.3),
             "■ 注目キーワード", size=12, color=GOOD, bold=True)
    # mine keywords (token-level, just simple counts of known terms)
    kw_pool = sum([kws for _, kws in CATEGORY_KEYWORDS], [])
    kw_counts: Counter = Counter()
    blob_all = " ".join((a.title + " " + a.summary) for a in rep.articles).lower()
    for kw in kw_pool:
        c = blob_all.count(kw.lower())
        if c > 0:
            kw_counts[kw] += c
    top_kw = kw_counts.most_common(18)
    # arrange as chips
    chip_x = R_X + Inches(0.25)
    chip_y = RB_Y + Inches(0.40)
    line_h = Inches(0.30)
    cur_x = chip_x
    cur_y = chip_y
    max_x = R_X + R_W - Inches(0.2)
    if top_kw:
        max_cnt = top_kw[0][1]
    for kw, cnt in top_kw:
        # size based on freq
        ratio = cnt / max_cnt if 'max_cnt' in locals() and max_cnt else 0.5
        size = 8 + int(ratio * 6)   # 8..14
        # estimate chip width: rough by char count
        text_w = Inches(0.16 + len(kw) * 0.10 + 0.35)
        if cur_x + text_w > max_x:
            cur_x = chip_x
            cur_y += line_h
            if cur_y + line_h > RB_Y + RB_H - Inches(0.1):
                break
        col = ACCENT if ratio > 0.66 else (WARN if ratio > 0.33 else COOL)
        add_rect(slide, cur_x, cur_y, text_w, Inches(0.26),
                 fill=PANEL_SOFT, line=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(slide, cur_x + Inches(0.05), cur_y + Inches(0.02),
                 text_w - Inches(0.1), Inches(0.22),
                 f"{kw} {cnt}", size=size, color=col, bold=True, align=PP_ALIGN.CENTER)
        cur_x += text_w + Inches(0.06)


def main():
    docs_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs")
    paths = sorted(glob.glob(os.path.join(docs_dir, "2026-*.md")))
    reports = [parse_md(p) for p in paths]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_index_slide(prs, reports)
    for r in reports:
        build_day_slide(prs, r)

    out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
    os.makedirs(out_dir, exist_ok=True)
    out = os.path.join(out_dir, "web3_news_digest_2026-05-07_2026-05-18.pptx")
    prs.save(out)
    print(f"Wrote {out}")
    print(f"Slides: {len(prs.slides)}  (1 cover + {len(reports)} daily)")
    for r in reports:
        print(f"  {r.date}  国内 {r.domestic_count:>3} / 海外 {r.overseas_count:>3}  "
              f"重大候補 {sum(1 for a in r.articles if a.hot_score >= 4)}")


if __name__ == "__main__":
    main()
